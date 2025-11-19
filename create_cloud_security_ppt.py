#!/usr/bin/env python3
"""
Cloud Security PowerPoint Presentation Generator
Creates a professional PowerPoint presentation about cloud security
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_cloud_security_presentation():
    """Create a comprehensive cloud security PowerPoint presentation"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    PRIMARY_COLOR = RGBColor(102, 126, 234)  # Blue-purple
    SECONDARY_COLOR = RGBColor(118, 75, 162)  # Purple
    ACCENT_COLOR = RGBColor(240, 147, 251)  # Pink
    TEXT_COLOR = RGBColor(30, 41, 59)  # Dark blue-gray
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background (simulated with shapes)
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Cloud Security: Protecting Your Digital Sky"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Simple Guide to Safeguarding Your Data in the Cloud"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "November 2025"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(255, 255, 255)
    date_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title = slide.shapes.title
    title.text = "Agenda"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Introduction to Cloud Security"
    
    agenda_items = [
        "Why It Matters",
        "Key Threats and Risks",
        "Best Practices",
        "Tools and Technologies",
        "Real-World Examples",
        "Future Trends",
        "Q&A"
    ]
    
    for item in agenda_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(24)
        p.space_before = Pt(12)
    
    # Slide 3: What is Cloud Security?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "What is Cloud Security?"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Definition"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "Cloud security protects data, applications, and infrastructure in cloud environments (AWS, Azure, Google Cloud) from threats."
    p.level = 0
    p.font.size = Pt(20)
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Analogy: The Cloud as a Digital Sky"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Think of the cloud as a vast digital sky—security is the 'air traffic control' ensuring safe flights for your data."
    p.level = 0
    p.font.size = Pt(20)
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Key Components"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    components = ["Identity and access management", "Data encryption", "Network security", "Compliance"]
    for comp in components:
        p = tf.add_paragraph()
        p.text = comp
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 4: Why Cloud Security Matters
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why Cloud Security Matters"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Benefits of Cloud"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "Scalability, cost savings, remote access—but introduces risks like data breaches"
    p.level = 0
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Critical Statistics"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "82% of organizations experienced a cloud security incident in 2023 (IBM report)"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Without security: fines, data loss, reputational damage"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Real Impact"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Losing customer data is like a storm wiping out your digital assets"
    p.level = 0
    p.font.size = Pt(20)
    p.font.italic = True
    
    # Slide 5: Key Threats and Risks
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Threats and Risks"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Common Threats"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    threats = [
        "Data breaches - unauthorized access to sensitive information",
        "DDoS attacks - overwhelming servers like a traffic jam",
        "Insider threats - employees or partners gone rogue",
        "Misconfigurations - accidental exposures (leaving doors unlocked)"
    ]
    
    for threat in threats:
        p = tf.add_paragraph()
        p.text = threat
        p.level = 1
        p.font.size = Pt(18)
        p.space_before = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "The CIA Triad"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    cia = ["Confidentiality", "Integrity", "Availability"]
    for item in cia:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 6: Best Practices (Part 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Best Practices - Part 1"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Secure Access"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "Use multi-factor authentication (MFA)"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Implement least-privilege access—only give keys to those who need them"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "2. Encrypt Data"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Protect data in transit and at rest"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Like locking valuables in a safe"
    p.level = 1
    p.font.size = Pt(20)
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "3. Regular Audits"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Monitor logs continuously"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Conduct penetration testing to find vulnerabilities"
    p.level = 1
    p.font.size = Pt(20)
    
    # Slide 7: Best Practices (Part 2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Best Practices - Part 2"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "4. Compliance and Policies"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "Follow standards like ISO 27001 or NIST"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Train employees on security hygiene"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "5. Incident Response"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Have a plan for breaches: Detect → Respond → Recover"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "6. Zero Trust Model"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = '"Never trust, always verify"'
    p.level = 1
    p.font.size = Pt(20)
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "Assume nothing is secure by default"
    p.level = 1
    p.font.size = Pt(20)
    
    # Slide 8: Tools and Technologies
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Tools and Technologies"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Popular Security Tools"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    tools = [
        "AWS GuardDuty - Threat detection",
        "Azure Security Center - Monitoring",
        "Cloudflare - DDoS protection",
        "Vault / AWS KMS - Encryption management"
    ]
    
    for tool in tools:
        p = tf.add_paragraph()
        p.text = tool
        p.level = 1
        p.font.size = Pt(20)
        p.space_before = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Emerging Technologies"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "AI-driven security for anomaly detection"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "💡 Tip: Start free with cloud provider tools; integrate with SIEM systems"
    p.level = 0
    p.font.size = Pt(18)
    p.font.italic = True
    p.space_before = Pt(20)
    
    # Slide 9: Real-World Examples
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Real-World Examples"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Success Story: Capital One (2019)"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "Breach led to better encryption practices"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Now uses advanced AI security systems"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Lesson Learned: SolarWinds (2020)"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Highlighted supply chain risks"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Emphasized importance of vendor vetting"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Key Takeaway"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Learn from incidents to strengthen your security posture"
    p.level = 0
    p.font.size = Pt(20)
    p.font.italic = True
    
    # Slide 10: Future Trends
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Trends in Cloud Security"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AI and Automation"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "Smarter threat prediction and response"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Edge Computing Security"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Protecting data closer to users"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Quantum-Resistant Encryption"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Preparing for quantum computing threats"
    p.level = 1
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "⚠️ Prediction: By 2025, 99% of cloud failures will be due to human error"
    p.level = 0
    p.font.size = Pt(18)
    p.font.italic = True
    p.space_before = Pt(20)
    
    # Slide 11: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Takeaway"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "Cloud security isn't optional—it's essential for a safe digital sky"
    p.level = 0
    p.font.size = Pt(22)
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "Call to Action"
    p.font.bold = True
    p.font.size = Pt(24)
    p.space_before = Pt(20)
    
    actions = [
        "Start small: Implement MFA today",
        "Conduct regular security audits",
        "Assess your cloud setup",
        "Contact security experts if needed"
    ]
    
    for action in actions:
        p = tf.add_paragraph()
        p.text = action
        p.level = 1
        p.font.size = Pt(20)
        p.space_before = Pt(8)
    
    p = tf.add_paragraph()
    p.text = '"Security is not a product, but a process." - Bruce Schneier'
    p.level = 0
    p.font.size = Pt(20)
    p.font.italic = True
    p.space_before = Pt(25)
    
    # Slide 12: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # Q&A Title
    qa_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1.5))
    qa_frame = qa_box.text_frame
    qa_frame.text = "Questions & Answers"
    qa_para = qa_frame.paragraphs[0]
    qa_para.font.size = Pt(54)
    qa_para.font.bold = True
    qa_para.font.color.rgb = RGBColor(255, 255, 255)
    qa_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(0.8))
    sub_frame = sub_box.text_frame
    sub_frame.text = "Thank you for your attention!"
    sub_para = sub_frame.paragraphs[0]
    sub_para.font.size = Pt(28)
    sub_para.font.color.rgb = RGBColor(255, 255, 255)
    sub_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "/vercel/sandbox/Cloud_Security_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ PowerPoint presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_cloud_security_presentation()
